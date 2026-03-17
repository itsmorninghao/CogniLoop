"""
Document parser — extract structured text from uploaded files.

Design principles from RAGFlow:
- Preserve document hierarchy (headings, sections, pages)
- Extract metadata (page numbers, section titles)
- Extensible via strategy pattern

Output format: list of ParsedSection with metadata for downstream chunking.
"""

from __future__ import annotations

import asyncio
import logging
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class ParsedSection:
    """A section of a parsed document, preserving structural context."""

    content: str
    page_number: int | None = None
    heading: str | None = None  # nearest heading above this content
    heading_level: int = 0  # 1=H1, 2=H2, etc.  0=body text
    section_path: str = ""  # e.g. "Chapter 1 > Section 2.1"
    metadata: dict = field(default_factory=dict)


@dataclass
class ParseResult:
    """Full result of parsing a document."""

    sections: list[ParsedSection]
    title: str = ""
    total_pages: int = 0
    metadata: dict = field(default_factory=dict)


class DocumentParser(ABC):
    """Abstract base class for document parsers — strategy pattern."""

    @abstractmethod
    def parse(self, file_path: str) -> ParseResult:
        """Parse a document and return structured sections."""
        ...

    @staticmethod
    def for_type(file_type: str) -> DocumentParser:
        """Factory method — return the correct parser for a file type."""
        file_type = file_type.upper()
        _parsers = {
            "PDF": PDFParser,
            "WORD": DocxParser,
            "DOCX": DocxParser,
            "PPT": PptxParser,
            "PPTX": PptxParser,
            "MARKDOWN": MarkdownParser,
            "MD": MarkdownParser,
            "TXT": PlainTextParser,
        }
        parser_cls = _parsers.get(file_type, PlainTextParser)
        return parser_cls()


class PDFParser(DocumentParser):
    """
    PDF parser powered by Docling AI vision models.
    Handles multi-column layouts, table extraction, scanned PDFs (OCR),
    and heading hierarchy detection.
    """

    def parse(self, file_path: str) -> ParseResult:
        from docling.document_converter import DocumentConverter
        from docling_core.types.doc import SectionHeaderItem, TableItem, TextItem

        converter = DocumentConverter()
        doc_result = converter.convert(file_path)
        doc = doc_result.document

        sections: list[ParsedSection] = []
        section_path_parts: list[str] = []
        current_heading = ""

        for item, _level in doc.iterate_items():
            if isinstance(item, SectionHeaderItem):
                heading_text = item.text.strip()
                heading_level = item.level  # Docling provides 1-6

                current_heading = heading_text
                while len(section_path_parts) >= heading_level:
                    section_path_parts.pop()
                section_path_parts.append(heading_text)

                sections.append(
                    ParsedSection(
                        content=heading_text,
                        heading_level=heading_level,
                        section_path=" > ".join(section_path_parts),
                        metadata={"type": "heading"},
                    )
                )

            elif isinstance(item, TableItem):
                table_md = item.export_to_markdown()
                if table_md.strip():
                    sections.append(
                        ParsedSection(
                            content=table_md,
                            heading=current_heading,
                            section_path=" > ".join(section_path_parts),
                            metadata={"type": "table"},
                        )
                    )

            elif isinstance(item, TextItem):
                text = item.text.strip()
                if text:
                    sections.append(
                        ParsedSection(
                            content=text,
                            heading=current_heading,
                            section_path=" > ".join(section_path_parts),
                            metadata={"type": "text"},
                        )
                    )

        title = doc.name or ""
        if not title:
            for s in sections:
                if s.heading_level == 1:
                    title = s.content
                    break

        return ParseResult(
            sections=sections,
            title=title,
            total_pages=len(doc.pages) if hasattr(doc, "pages") else 0,
            metadata={"parser": "DoclingPDFParser", "file_path": file_path},
        )


class DocxParser(DocumentParser):
    """Word parser — leverages native paragraph styles for heading detection."""

    def parse(self, file_path: str) -> ParseResult:
        from docx import Document

        doc = Document(file_path)
        sections: list[ParsedSection] = []
        current_heading = ""
        section_path_parts: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            heading_level = self._style_to_level(para.style.name)

            if heading_level > 0:
                current_heading = text
                while len(section_path_parts) >= heading_level:
                    section_path_parts.pop()
                section_path_parts.append(current_heading)

            sections.append(
                ParsedSection(
                    content=text,
                    heading=current_heading if heading_level == 0 else None,
                    heading_level=heading_level,
                    section_path=" > ".join(section_path_parts),
                    metadata={"style": para.style.name},
                )
            )

        title = ""
        for s in sections:
            if s.heading_level == 1:
                title = s.content
                break

        return ParseResult(
            sections=sections,
            title=title,
            metadata={"parser": "DocxParser"},
        )

    @staticmethod
    def _style_to_level(style_name: str) -> int:
        style_lower = style_name.lower()
        if "heading 1" in style_lower or "标题 1" in style_lower:
            return 1
        if "heading 2" in style_lower or "标题 2" in style_lower:
            return 2
        if "heading 3" in style_lower or "标题 3" in style_lower:
            return 3
        if "heading" in style_lower or "标题" in style_lower:
            return 2
        if style_lower == "title":
            return 1
        return 0


class PptxParser(DocumentParser):
    """PowerPoint parser — one section per slide with slide title as heading."""

    def parse(self, file_path: str) -> ParseResult:
        from pptx import Presentation

        prs = Presentation(file_path)
        sections: list[ParsedSection] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_title = ""
            body_texts = []

            for shape in slide.shapes:
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue
                if shape.shape_type is not None and "TITLE" in str(shape.shape_type):
                    slide_title = shape.text.strip()
                else:
                    body_texts.append(shape.text.strip())

            if not slide_title and body_texts:
                slide_title = body_texts[0][:50]

            content = "\n".join(body_texts) if body_texts else slide_title
            if content.strip():
                sections.append(
                    ParsedSection(
                        content=content,
                        page_number=slide_idx + 1,
                        heading=slide_title,
                        heading_level=2,
                        section_path=f"Slide {slide_idx + 1}: {slide_title}",
                        metadata={"slide_index": slide_idx + 1},
                    )
                )

        return ParseResult(
            sections=sections,
            title=sections[0].heading if sections else "",
            total_pages=len(prs.slides),
            metadata={"parser": "PptxParser"},
        )


class MarkdownParser(DocumentParser):
    """Markdown parser — uses # heading syntax for structure."""

    def parse(self, file_path: str) -> ParseResult:
        with open(file_path, encoding="utf-8", errors="ignore") as f:
            content = f.read()

        sections: list[ParsedSection] = []
        current_heading = ""
        section_path_parts: list[str] = []
        buffer = []

        for line in content.split("\n"):
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading_match:
                if buffer:
                    body = "\n".join(buffer).strip()
                    if body:
                        sections.append(
                            ParsedSection(
                                content=body,
                                heading=current_heading,
                                section_path=" > ".join(section_path_parts),
                            )
                        )
                    buffer = []

                level = len(heading_match.group(1))
                heading_text = heading_match.group(2).strip()
                current_heading = heading_text

                while len(section_path_parts) >= level:
                    section_path_parts.pop()
                section_path_parts.append(heading_text)

                sections.append(
                    ParsedSection(
                        content=heading_text,
                        heading_level=level,
                        section_path=" > ".join(section_path_parts),
                    )
                )
            else:
                buffer.append(line)

        if buffer:
            body = "\n".join(buffer).strip()
            if body:
                sections.append(
                    ParsedSection(
                        content=body,
                        heading=current_heading,
                        section_path=" > ".join(section_path_parts),
                    )
                )

        title = ""
        for s in sections:
            if s.heading_level > 0:
                title = s.content
                break

        return ParseResult(
            sections=sections,
            title=title,
            metadata={"parser": "MarkdownParser"},
        )


class PlainTextParser(DocumentParser):
    """Fallback parser for plain text files."""

    def parse(self, file_path: str) -> ParseResult:
        with open(file_path, encoding="utf-8", errors="ignore") as f:
            content = f.read()

        paragraphs = [p.strip() for p in re.split(r"\n\s*\n", content) if p.strip()]

        sections = [
            ParsedSection(content=p, metadata={"paragraph_index": i})
            for i, p in enumerate(paragraphs)
        ]

        return ParseResult(
            sections=sections,
            metadata={"parser": "PlainTextParser"},
        )


async def parse_document(file_path: str, file_type: str) -> ParseResult:
    """Parse a document and return structured sections."""
    parser = DocumentParser.for_type(file_type)
    result = await asyncio.to_thread(parser.parse, file_path)
    logger.info(
        "Parsed %s (%s): %d sections, title='%s'",
        file_path,
        file_type,
        len(result.sections),
        result.title[:60],
    )
    return result
